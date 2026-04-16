"""Generate a 15-minute pitch PowerPoint for The Arc Mercer proposal."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=DARK_BLUE):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(tf, items, font_size=16, color=DARK_GRAY, bold_prefix=False, spacing=Pt(6)):
    """Add bullet items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.font.size = Pt(font_size)

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = "• " + prefix + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = "Calibri"
        else:
            p.text = "• " + item
            p.font.color.rgb = color
            p.font.name = "Calibri"


def section_header_bar(slide, text):
    """Add a colored header bar at the top of a content slide."""
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                text, font_size=32, color=WHITE, bold=True)


def add_picture(slide, filename, left, top, width):
    """Add an image to a slide; height is auto-calculated to preserve aspect ratio."""
    img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), filename)
    return slide.shapes.add_picture(img_path, left, top, width=width)


# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(3.2), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.6),
            "IoTeye Inc.", font_size=24, color=LIGHT_BLUE, bold=False)
add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(1.2),
            "Employee Engagement &\nRecognition Platform", font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "Proposal for The Arc Mercer", font_size=24, color=WHITE, bold=False)
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "April 2026  |  Basking Ridge, NJ  |  www.ioteyeinc.com", font_size=16, color=LIGHT_BLUE)


# ─── SLIDE 2: Agenda ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Agenda")

items = [
    "1.  About IoTeye & Platform Ecosystem",
    "2.  BrainClaw AI Agents \u2014 Already at Arc Mercer",
    "3.  Understanding Your Challenge",
    "4.  Platform Overview & Architecture",
    "5.  Key Features (Nominations \u2192 Engagement \u2192 Marketplace)",
    "6.  Integration, Authentication & Accessibility",
    "7.  Dashboards \u2014 Employee & Leadership",
    "8.  Implementation Timeline & Pricing",
    "9.  Feature Comparison \u2014 How We Stack Up",
    "10. Why IoTeye \u2014 Next Steps",
]
tf = add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5),
                 "", font_size=22, color=DARK_GRAY)
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    p.space_after = Pt(10)


# ─── SLIDE 3: About IoTeye ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "About IoTeye Inc.")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "Who We Are", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                 "", font_size=16, color=DARK_GRAY)
add_bullet_list(tf, [
    "Remote software company — Basking Ridge, NJ",
    "5+ years building special care software",
    "1,000+ consumers across our platforms",
    "Core domain: Special Care sector",
    "Full-stack SaaS expertise: design → deploy → support",
], font_size=16, spacing=Pt(10))

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
            "Our Products", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(7), Inches(2.0), Inches(5.8), Inches(4.8),
                 "", font_size=15, color=DARK_GRAY)
add_bullet_list(tf, [
    "SpringBoard: Multi-agency route management & operations (route opt., fleet, Twilio SMS, multi-tenant)",
    "BrainBook: AI platform \u2014 Next.js/React/TypeScript, Supabase, AI agents, voice I/O, desktop & web",
    "Guardian: Notification system \u2014 SMS, push notifications, reply tracking, iOS & Android",
    "BrainClaw: AI agent gateway (GPT-4.1) \u2014 37-tool SpringBoard agent & 57-tool Samsara agent live at Arc Mercer",
], font_size=14, bold_prefix=False, spacing=Pt(10))


# \u2500\u2500\u2500 SLIDE 4: Total Agency Services for Special Care \u2500\u2500\u2500
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Total Agency Services for Special Care")

# Left panel
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.5), LIGHT_GRAY)
add_textbox(slide, Inches(0.65), Inches(1.4), Inches(5.3), Inches(0.45),
            "Special Care Cloud (AWS)", font_size=18, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.65), Inches(1.95), Inches(5.3), Inches(4.8),
                 "", font_size=13, color=DARK_GRAY)
add_bullet_list(tf, [
    "Market: 6,000+ special care agencies in the USA",
    "Agentic AI Platform \u2014 Cloud, Desktop, Mobile, IoT",
    "Personal AI assistants for agency workflows",
    "Apps for Parent/Guardian, Consumer/Patient, Staff/Driver",
    "IoT: OBD, DashCam, Minew Proximity, Home Assistant Green",
    "Samsara App Partner \u2014 fleet telematics",
    "HIPAA-compliant AI inference servers",
], font_size=13, color=DARK_GRAY, spacing=Pt(8))

# Right panel
add_shape_bg(slide, Inches(7.0), Inches(1.3), Inches(5.9), Inches(5.5), LIGHT_GRAY)
add_textbox(slide, Inches(7.25), Inches(1.4), Inches(5.4), Inches(0.45),
            "Platform Products", font_size=18, color=ACCENT_BLUE, bold=True)
products_info = [
    ("SpringBoard", "Route/Fleet \u2014 paratransit ops, multi-tenant"),
    ("Guardian", "Notifications \u2014 SMS, push, reply tracking"),
    ("BrainBook", "AI Platform \u2014 desktop, web, mobile, voice I/O"),
    ("BrainClaw", "AI Agents \u2014 94+ tools, GPT-4.1, Arc Mercer live"),
]
y_p = Inches(2.05)
for pname, pdesc in products_info:
    add_textbox(slide, Inches(7.25), y_p, Inches(5.4), Inches(0.3),
                pname, font_size=15, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(7.25), y_p + Inches(0.3), Inches(5.4), Inches(0.35),
                pdesc, font_size=13, color=DARK_GRAY)
    y_p += Inches(0.85)


# \u2500\u2500\u2500 SLIDE 5: BrainClaw AI Agents \u2500\u2500\u2500
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "BrainClaw AI Agents \u2014 Already at Arc Mercer")

# Left column: In production today
add_textbox(slide, Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.45),
            "In Production at Arc Mercer Today", font_size=20, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.35),
            "Powered by BrainClaw + GPT-4.1 (GitHub Copilot)", font_size=14, color=MED_GRAY)
tf = add_textbox(slide, Inches(0.5), Inches(2.25), Inches(6.0), Inches(2.1),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "SpringBoard Agent \u2014 37 tools",
    "  Routes, fleet, consumers, optimization, SMS, GPS, SOS",
    "Samsara Agent \u2014 57 tools",
    "  Vehicles, drivers, HOS, safety events, DVIRs, GPS",
], font_size=14, spacing=Pt(5))

add_textbox(slide, Inches(0.5), Inches(4.45), Inches(6.0), Inches(0.4),
            "Arc Mercer staff use natural language today:", font_size=14, color=DARK_GRAY, bold=True)
add_shape_bg(slide, Inches(0.5), Inches(4.9), Inches(6.0), Inches(2.0), LIGHT_GRAY)
add_textbox(slide, Inches(0.65), Inches(4.95), Inches(5.7), Inches(1.8),
            '"Optimize routes for Weekday AM"\n"Where is vehicle 147 right now?"\n"SMS all Residential drivers \u2014 10 min delay"\n"Generate daily manifest for Arc Mercer"',
            font_size=13, color=DARK_GRAY, font_name="Consolas")

# Right column: Extension to Recognition Platform
add_shape_bg(slide, Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.8), LIGHT_BLUE)
add_textbox(slide, Inches(7.1), Inches(1.3), Inches(5.7), Inches(0.45),
            "Extending to Recognition Platform", font_size=20, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(7.1), Inches(1.85), Inches(5.7), Inches(0.4),
            "Natural language for HR & leadership:", font_size=15, color=DARK_GRAY, bold=True)
tf = add_textbox(slide, Inches(7.1), Inches(2.35), Inches(5.7), Inches(3.2),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    '"Which program had most nominations last month?"',
    '"Who hasn\'t been recognized in 60+ days?"',
    '"Flag unusual point bestowals this week"',
    '"Generate Q1 engagement report as CSV"',
    '"What\'s our engagement score this month?"',
], font_size=14, spacing=Pt(10))

add_shape_bg(slide, Inches(7.0), Inches(5.8), Inches(5.8), Inches(0.8), ACCENT_BLUE)
add_textbox(slide, Inches(7.2), Inches(5.85), Inches(5.4), Inches(0.65),
            "\u2605 Optional add-on \u2014 same BrainBook UI Arc Mercer already uses",
            font_size=14, color=WHITE, bold=True)


# \u2500\u2500\u2500 SLIDE 6: Understanding the Challenge \u2500\u2500\u2500
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Understanding Your Challenge")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "The Arc Mercer's Current State", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4),
                 "", font_size=16, color=DARK_GRAY)
add_bullet_list(tf, [
    "Manual recognition processes — inconsistent, hard to track",
    "Multiple programs (Residential, Day Programs) operating in silos",
    "Frontline staff with varying technical literacy",
    "Many frontline DSPs lack company email addresses",
    "Selection committee uses ranked-choice voting — no digital tool",
    "No centralized system for nominations or point tracking",
    "No connection to Paycom HRIS for employee data",
], font_size=16, spacing=Pt(10))

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
            "What You Need", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4),
                 "", font_size=16, color=DARK_GRAY)
add_bullet_list(tf, [
    "Centralized, web-based recognition platform",
    "Program-based competition engine with committee voting",
    "Multi-stream point system with audit trails",
    "Rewards marketplace — digital, caf\u00e9, and in-person store",
    "Paycom integration & flexible auth (SSO + Badge ID + PIN)",
    "Community engagement: TAC, culture groups, volunteering",
    "Mobile-first, accessible to all staff (incl. dark mode)",
], font_size=16, spacing=Pt(10))


# ─── SLIDE 5: Solution Overview ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Our Solution — Platform Overview")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "Custom-Built SaaS — Tailored for The Arc Mercer", font_size=22, color=ACCENT_BLUE, bold=True)

# Architecture layers
layers = [
    ("React + TypeScript PWA", "Mobile-first responsive design, offline capable", ACCENT_BLUE),
    ("Supabase (PostgreSQL)", "Auth, real-time subscriptions, Row Level Security", RGBColor(0x1A, 0x5C, 0x8A)),
    ("Python (FastAPI)", "Paycom SFTP sync, gift card APIs, scheduled jobs", RGBColor(0x14, 0x73, 0x5A)),
    ("AWS Cloud", "ECS Fargate, CloudFront, S3, WAF, CloudWatch", RGBColor(0xE6, 0x7E, 0x22)),
]

y_pos = Inches(2.0)
for label, desc, color in layers:
    shape = add_shape_bg(slide, Inches(1.5), y_pos, Inches(10), Inches(1.0), color)
    add_textbox(slide, Inches(2), y_pos + Inches(0.08), Inches(4), Inches(0.45),
                label, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, Inches(2), y_pos + Inches(0.48), Inches(9), Inches(0.45),
                desc, font_size=14, color=WHITE)
    y_pos += Inches(1.15)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            "Security: AES-256 at rest  |  TLS 1.3 in transit  |  WCAG AAA  |  Full audit logging",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 6: Nomination Engine ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Nominations, Voting & QA Engine")

# Left column: Nominations
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.45),
            "Nominations", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(1.85), Inches(5.5), Inches(1.6),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "3 award types: Employee of the Month, Rising Star, Team Impact",
    "Step-by-step guided form with large buttons",
    "1 nomination per candidate/month, 5 total/month",
], font_size=14, spacing=Pt(6))

# Left column: Selection Committee Voting
add_textbox(slide, Inches(0.4), Inches(3.4), Inches(5.5), Inches(0.45),
            "Selection Committee Voting", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(3.9), Inches(5.5), Inches(1.5),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Ranked-choice drag-and-drop voting interface",
    "Private ballots \u2014 not visible until window closes",
    "Automated tally (Borda count / weighted rank)",
    "Per-voter audit trail & historical ballot archive",
], font_size=14, spacing=Pt(6))

# Left column: QA Queue + Revise-and-Resubmit
add_textbox(slide, Inches(0.4), Inches(5.4), Inches(5.5), Inches(0.45),
            "QA Queue & Revise-Resubmit", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(5.9), Inches(5.5), Inches(1.4),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Approve / Return for Edits / Decline with batch operations",
    "Formal resubmission loop: status tracking through full cycle",
    "Admin comments \u2192 nominator notification \u2192 edit & resubmit",
], font_size=14, spacing=Pt(6))

# Right: screenshot of Nomination Center
add_picture(slide, "reward_2.png", Inches(6.1), Inches(1.5), Inches(6.8))


# ─── SLIDE 7: Points & Engagement ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Points, Engagement & Community Modules")

# Point streams table
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
            "Multi-Stream Point Accrual", font_size=22, color=ACCENT_BLUE, bold=True)

streams = [
    ("Point Stream", "Trigger"),
    ("Employee / Rising Star / Team Impact (Approved)", "QA queue approval"),
    ("TAC Meeting Attendance", "Admin check-in"),
    ("TAC Meeting Role Completion", "Role fulfilled"),
    ("Culture Group Participation", "Leader/admin log"),
    ("Volunteer / Event Help", "Post-event confirmation"),
    ("Social Media Engagement", "Screenshot review"),
    ("Newsletter Trivia", "Correct answer"),
    ("Birthday / Anniversary", "Automated (Paycom)"),
    ("Leadership Bestowal", "Manual award"),
]

y = Inches(2.0)
for i, (stream, trigger) in enumerate(streams):
    bg_color = ACCENT_BLUE if i == 0 else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    txt_color = WHITE if i == 0 else DARK_GRAY
    add_shape_bg(slide, Inches(0.8), y, Inches(7), Inches(0.38), bg_color)
    add_textbox(slide, Inches(0.9), y + Inches(0.01), Inches(4.5), Inches(0.35),
                stream, font_size=12, color=txt_color, bold=(i == 0))
    add_textbox(slide, Inches(5.5), y + Inches(0.01), Inches(2.2), Inches(0.35),
                trigger, font_size=12, color=txt_color, bold=(i == 0))
    y += Inches(0.38)

add_textbox(slide, Inches(0.8), y + Inches(0.15), Inches(7), Inches(0.4),
            "All point values are admin-configurable \u2014 no code changes needed",
            font_size=14, color=GREEN, bold=True)

# Community Engagement Modules
add_textbox(slide, Inches(8.5), Inches(1.4), Inches(4.5), Inches(0.5),
            "Community Engagement", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(8.5), Inches(2.0), Inches(4.5), Inches(5),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "TAC RSVP & Role Sign-Ups: RSVP workflow, role selection (setup, check-in, cleanup), automated reminders",
    "Culture Groups: Self-service directory, one-tap sign-up, attendance tracking, participation points",
    "Volunteer Forms: Event help board, volunteer slots, completion tracking, engagement score",
    "Social Media: Screenshot upload, admin review queue, campaign tracking",
    "Audit & Transparency: Immutable transaction log, anomaly detection, CSV export",
], font_size=13, bold_prefix=True, spacing=Pt(10))


# ─── SLIDE 8: Rewards Marketplace ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Rewards Marketplace \u2014 Digital, Caf\u00e9 & In-Person Store")

# Left column: Digital Storefront
add_textbox(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.45),
            "Digital Storefront", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(1.85), Inches(5.5), Inches(1.5),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Visual catalog with images, point costs, descriptions",
    "Categories: Merchandise, Cafeteria Meals, Gift Cards",
    "Real-time point balance during browsing",
    "One-tap redemption with confirmation",
], font_size=14, spacing=Pt(6))

# Left column: Cafe Cashier Redemption
add_textbox(slide, Inches(0.4), Inches(3.35), Inches(5.5), Inches(0.45),
            "Caf\u00e9 Cashier Redemption Screen", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(3.85), Inches(5.5), Inches(1.3),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Tablet-friendly cashier interface for meal redemptions",
    "Employee lookup by name, badge ID, or QR code",
    "Balance display \u2192 meal selection \u2192 point deduction",
    "No POS integration required \u2014 standalone operation",
], font_size=14, spacing=Pt(6))

# Left column: In-Person Store QR
add_textbox(slide, Inches(0.4), Inches(5.15), Inches(5.5), Inches(0.45),
            "In-Person Store & QR Redemption", font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.4), Inches(5.65), Inches(5.5), Inches(1.3),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Employee QR code on profile / home screen",
    "Store attendant scans QR \u2192 select item \u2192 confirm",
    "Inventory sync: physical store + digital storefront",
], font_size=14, spacing=Pt(6))

# Right: screenshot of Rewards Marketplace
add_picture(slide, "reward_4.png", Inches(6.1), Inches(1.5), Inches(6.8))


# ─── SLIDE 9: Paycom & SSO ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Paycom Integration & Authentication")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "Paycom HRIS Integration", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
                 "", font_size=15, color=DARK_GRAY)
add_bullet_list(tf, [
    "Primary: SFTP nightly file export (CSV/XML)",
    "Python ETL service parses & syncs to Supabase",
    "Auto-provision new hires, deactivate terminations",
    "Sync: Employee ID, name, email, program, hire date, DOB",
    "Discrepancy reports flagged for HR review",
    "Optional: REST API enhancement (if subscription supports)",
], font_size=15, spacing=Pt(8))

# Architecture mini-diagram as text
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.2), LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(4.9), Inches(5.2), Inches(2.0),
            "Paycom \u2192 SFTP Export \u2192 AWS S3 \u2192 Python ETL\n          \u2193\n    Supabase (PostgreSQL)\n          \u2193\n    React Front End (Live Data)",
            font_size=13, color=DARK_GRAY, font_name="Consolas")

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
            "Three Login Paths", font_size=24, color=ACCENT_BLUE, bold=True)

# Path 1: Badge ID + PIN
add_shape_bg(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(1.35), LIGHT_BLUE)
add_textbox(slide, Inches(7.15), Inches(2.05), Inches(5.2), Inches(0.35),
            "\u2780  Badge ID + PIN (Frontline Staff)", font_size=16, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(7.15), Inches(2.4), Inches(5.2), Inches(0.9),
                 "", font_size=13, color=DARK_GRAY)
add_bullet_list(tf, [
    "Paycom badge ID \u2014 the ID they already carry daily",
    "Self-set 4\u20136 digit PIN, no email required",
    "Hashed/salted PINs, lockout after 5 failed attempts",
], font_size=13, spacing=Pt(4))

# Path 2: Google SSO
add_shape_bg(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(1.2), LIGHT_GRAY)
add_textbox(slide, Inches(7.15), Inches(3.55), Inches(5.2), Inches(0.35),
            "\u2781  Google Workspace SSO (Admin & Mgmt)", font_size=16, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(7.15), Inches(3.9), Inches(5.2), Inches(0.7),
                 "", font_size=13, color=DARK_GRAY)
add_bullet_list(tf, [
    "Supabase Auth: SAML 2.0, OAuth 2.0, OIDC",
    "Future: Microsoft Entra ID \u2014 zero-downtime migration",
], font_size=13, spacing=Pt(4))

# Path 3: Email/Password
add_shape_bg(slide, Inches(7), Inches(4.85), Inches(5.5), Inches(0.9), LIGHT_GRAY)
add_textbox(slide, Inches(7.15), Inches(4.9), Inches(5.2), Inches(0.35),
            "\u2782  Email / Password Fallback", font_size=16, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(7.15), Inches(5.25), Inches(5.2), Inches(0.5),
                 "", font_size=13, color=DARK_GRAY)
add_bullet_list(tf, [
    "For part-time or seasonal staff with email but not on SSO",
    "Guardian authentication",
], font_size=13, spacing=Pt(4))

# Key callout
add_shape_bg(slide, Inches(7), Inches(5.8), Inches(5.5), Inches(0.8), GREEN)
add_textbox(slide, Inches(7.15), Inches(5.85), Inches(5.2), Inches(0.65),
            "\u2605 Every employee can access the platform \u2014\nno matter their tech level or email situation",
            font_size=14, color=WHITE, bold=True)


# ─── SLIDE 10: Dashboards ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Dashboards — Employee & Leadership")

# Column labels
add_textbox(slide, Inches(0.15), Inches(1.15), Inches(6.0), Inches(0.38),
            "Employee Dashboard", font_size=18, color=ACCENT_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.15), Inches(1.15), Inches(6.0), Inches(0.38),
            "Leadership Console", font_size=18, color=ACCENT_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)

# Screenshots side by side (1456×816 px → ratio 1.784; width 6.0" → height ≈3.36")
add_picture(slide, "reward_1.png", Inches(0.15), Inches(1.55), Inches(6.0))
add_picture(slide, "reward_3.png", Inches(7.15), Inches(1.55), Inches(6.0))

# Captions below screenshots (bottom ≈ 1.55 + 3.36 = 4.91")
add_textbox(slide, Inches(0.15), Inches(5.05), Inches(6.0), Inches(0.35),
            "Points Balance  •  Wall of Fame  •  Quick Nominate",
            font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.15), Inches(5.05), Inches(6.0), Inches(0.35),
            "Total Points  •  Engagement Rate  •  QA Pending Review",
            font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom summary bar
add_shape_bg(slide, Inches(0), Inches(5.55), SLIDE_W, Inches(0.7), LIGHT_BLUE)
add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12), Inches(0.6),
            "All dashboards are real-time  •  Mobile-first  •  Accessible on any device  •  Role-based visibility",
            font_size=15, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 11: Phased Timeline ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Phased Implementation — 36 Weeks")

phases = [
    ("Phase 1: Foundation", "Weeks 1\u201312", "Nominations, committee voting, Paycom sync, Badge ID + PIN auth, SSO, QA queue",
     ACCENT_BLUE, Inches(1.5), Inches(4.8)),
    ("Phase 2: Engagement", "Weeks 13\u201324", "Dashboards, TAC RSVP, culture groups, volunteer/social media, dark mode, notifications",
     RGBColor(0x14, 0x73, 0x5A), Inches(5.3), Inches(3.5)),
    ("Phase 3: Marketplace", "Weeks 25\u201336", "Digital store, caf\u00e9 cashier screen, in-person QR store, gift cards, fulfillment",
     ORANGE, Inches(8.8), Inches(3.5)),
]

y_base = Inches(2.0)
for label, weeks, desc, color, x_start, width in phases:
    # Timeline bar
    add_shape_bg(slide, x_start, y_base, width, Inches(0.6), color)
    add_textbox(slide, x_start + Inches(0.2), y_base + Inches(0.05), width - Inches(0.3), Inches(0.5),
                f"{label}  ({weeks})", font_size=16, color=WHITE, bold=True)
    y_base += Inches(0.8)

# Details below
y_detail = Inches(4.8)
details = [
    ("Phase 1 Deliverables", [
        "Nomination engine + ranked-choice committee voting",
        "QA queue with revise-and-resubmit workflow",
        "Paycom SFTP sync + audit trail",
        "Auth: Badge ID + PIN, Google SSO, email fallback",
    ], ACCENT_BLUE),
    ("Phase 2 Deliverables", [
        "Employee & leadership dashboards",
        "Dark mode / light mode toggle",
        "TAC RSVP, culture groups, volunteer forms",
        "Social media engagement (screenshot)",
        "Push/email notifications",
    ], RGBColor(0x14, 0x73, 0x5A)),
    ("Phase 3 Deliverables", [
        "Digital rewards storefront",
        "Caf\u00e9 cashier redemption screen",
        "In-person store QR code scanning",
        "Gift card instant delivery",
        "Inventory management & fulfillment",
    ], ORANGE),
]

x_pos = Inches(0.8)
for title, items, color in details:
    add_textbox(slide, x_pos, y_detail, Inches(3.8), Inches(0.4),
                title, font_size=16, color=color, bold=True)
    tf = add_textbox(slide, x_pos, y_detail + Inches(0.4), Inches(3.8), Inches(2.5),
                     "", font_size=13, color=DARK_GRAY)
    add_bullet_list(tf, items, font_size=13, spacing=Pt(4))
    x_pos += Inches(4.2)


# ─── SLIDE 12: Pricing — Option A ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Pricing — Option A: Source-Code Ownership")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "The Arc Mercer owns all source code, database, and IP",
            font_size=18, color=GREEN, bold=True)

# Implementation cost table
headers = [("Phase", Inches(4)), ("Hours", Inches(1.5)), ("Cost", Inches(1.5))]
rows_data = [
    ("Phase 1: Nominations, Voting, Paycom, Badge ID Auth, SSO", "~480 hrs", "$45,000"),
    ("Phase 2: Dashboards, Engagement Modules, Dark Mode, Notifications", "~310 hrs", "$35,000"),
    ("Phase 3: Marketplace, Caf\u00e9 Redemption, QR Store, Gift Cards", "~280 hrs", "$30,000"),
    ("Total Implementation", "~1,070 hrs", "$110,000"),
]

y = Inches(1.9)
x_base = Inches(1.5)
# Header row
for j, (h, w) in enumerate(headers):
    x = x_base + sum(hw[1] for hw in headers[:j])
    add_shape_bg(slide, x, y, w, Inches(0.45), ACCENT_BLUE)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.4),
                h, font_size=14, color=WHITE, bold=True)
y += Inches(0.45)

for i, (phase, hours, cost) in enumerate(rows_data):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    is_total = (i == len(rows_data) - 1)
    for j, (val, w) in enumerate(zip([phase, hours, cost], [h[1] for h in headers])):
        x = x_base + sum(hw[1] for hw in headers[:j])
        add_shape_bg(slide, x, y, w, Inches(0.45), DARK_BLUE if is_total else bg)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.4),
                    val, font_size=13, color=WHITE if is_total else DARK_GRAY, bold=is_total)
    y += Inches(0.45)

add_textbox(slide, Inches(1.5), y + Inches(0.1), Inches(7), Inches(0.4),
            "Blended rate: ~$103/hr — well below $150-250/hr industry standard",
            font_size=14, color=GREEN, bold=True)

# Annual costs
add_textbox(slide, Inches(8.5), Inches(1.9), Inches(4.5), Inches(0.5),
            "Annual Recurring", font_size=20, color=ACCENT_BLUE, bold=True)
annual_items = [
    "AWS hosting/CDN/monitoring: $6,000",
    "Supabase Pro Plan: $3,000",
    "SSL & domain: $200",
    "Maintenance & support: $12,000",
    "Total annual: $21,200/year",
]
tf = add_textbox(slide, Inches(8.5), Inches(2.5), Inches(4.5), Inches(3),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, annual_items, font_size=14, spacing=Pt(8))

add_shape_bg(slide, Inches(8.5), Inches(5.0), Inches(4), Inches(1.2), LIGHT_BLUE)
add_textbox(slide, Inches(8.7), Inches(5.1), Inches(3.6), Inches(0.4),
            "Year 1 Total: $131,200", font_size=18, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(8.7), Inches(5.5), Inches(3.6), Inches(0.4),
            "Year 2+: $21,200/year", font_size=18, color=DARK_BLUE, bold=True)

# Add-ons callout banner
add_shape_bg(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.7), GREEN)
add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.6),
            "\u2605  AI-assisted add-ons INCLUDED \u2014 LLM usage fee only (pass-through)  |  Major redesigns not covered  \u2605",
            font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 13: Pricing — Option B & Option C ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Pricing \u2014 Option B (SaaS) & Option C (Shared Ownership)")

# Option B (left)
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(3.0), LIGHT_BLUE)
add_textbox(slide, Inches(0.6), Inches(1.35), Inches(5.5), Inches(0.45),
            "Option B: IoTeye-Hosted SaaS", font_size=20, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.6), Inches(1.85), Inches(5.5), Inches(2.3),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Implementation fee: $0",
    "Subscription: $2,500/month ($30,000/year)",
    "3-year minimum commitment",
    "IoTeye retains all source code ownership",
    "Hosting, maintenance, updates all included",
], font_size=14, spacing=Pt(7))

# Option C (right)
add_shape_bg(slide, Inches(7.0), Inches(1.3), Inches(5.9), Inches(3.0), LIGHT_GRAY)
add_textbox(slide, Inches(7.2), Inches(1.35), Inches(5.5), Inches(0.45),
            "Option C: Shared Ownership (Hybrid)", font_size=20, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(7.2), Inches(1.85), Inches(5.5), Inches(2.3),
                 "", font_size=14, color=DARK_GRAY)
add_bullet_list(tf, [
    "Implementation fee: $55,000 (50% of full build)",
    "Annual platform fee: $30,000/year (managed service)",
    "Arc Mercer & IoTeye share recognition platform code (50/50)",
    "Excludes: BrainBook, SpringBoard, Guardian, BrainClaw",
    "Reuse rights deferred to agreement",
], font_size=14, spacing=Pt(7))

# Add-ons promo badge
add_shape_bg(slide, Inches(0.4), Inches(4.4), Inches(12.5), Inches(0.45), GREEN)
add_textbox(slide, Inches(0.6), Inches(4.43), Inches(12), Inches(0.4),
            "\u2605  AI-assisted add-ons INCLUDED \u2014 LLM usage fee only (pass-through)  |  Major redesigns not covered  \u2605",
            font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 3-column comparison table
add_textbox(slide, Inches(0.4), Inches(5.0), Inches(12), Inches(0.4),
            "Option A vs. B vs. C \u2014 Cost Comparison", font_size=18, color=DARK_BLUE, bold=True)

comp_headers = [("", Inches(2.2)), ("A: Own Code", Inches(2.6)), ("B: SaaS", Inches(2.6)), ("C: Shared", Inches(2.6))]
comp_rows = [
    ("Upfront Cost", "$110,000", "$0", "$55,000"),
    ("Annual Cost", "$21,200/yr", "$30,000/yr", "$30,000/yr"),
    ("3-Year Total", "$173,600", "$90,000", "$145,000"),
    ("5-Year Total", "$216,000", "$150,000", "$205,000"),
    ("Code Ownership", "Full", "None", "50/50 recognition"),
]

y = Inches(5.4)
x_base = Inches(1.2)
for j, (h, w) in enumerate(comp_headers):
    x = x_base + sum(hw[1] for hw in comp_headers[:j])
    add_shape_bg(slide, x, y, w, Inches(0.35), ACCENT_BLUE)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.01), w - Inches(0.16), Inches(0.32),
                h, font_size=12, color=WHITE, bold=True)
y += Inches(0.35)

for i, (label, a, b, c) in enumerate(comp_rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (val, w) in enumerate(zip([label, a, b, c], [h[1] for h in comp_headers])):
        x = x_base + sum(hw[1] for hw in comp_headers[:j])
        add_shape_bg(slide, x, y, w, Inches(0.32), bg)
        add_textbox(slide, x + Inches(0.08), y + Inches(0.01), w - Inches(0.16), Inches(0.28),
                    val, font_size=11, color=DARK_GRAY, bold=(j == 0))
    y += Inches(0.32)


# ─── SLIDE 14: Feature Comparison Summary ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Feature Comparison \u2014 How We Stack Up")

add_textbox(slide, Inches(0.8), Inches(1.25), Inches(8), Inches(0.45),
            "37 features evaluated across all 3 proposals submitted for this RFP",
            font_size=16, color=MED_GRAY)

# Category scorecard table
cat_headers = [
    ("Category", Inches(3.8)),
    ("IoTeye", Inches(1.2)),
    ("Proposal 1", Inches(1.4)),
    ("Proposal 2", Inches(1.4)),
]
cat_rows = [
    ("Nomination & Recognition (6)", "6", "5", "4"),
    ("Selection & Voting (3)", "3", "2", "3"),
    ("Authentication & Access (5)", "4", "1", "1"),
    ("Engagement & Participation (6)", "6", "5", "2"),
    ("Rewards & Redemption (5)", "5", "4", "4"),
    ("Analytics, AI & Reporting (4)", "4", "2", "2"),
    ("Technical & Infrastructure (8)", "8", "2", "2"),
    ("TOTAL", "36 / 37", "21 / 37", "18 / 37"),
]

y_t = Inches(1.8)
x_base_t = Inches(0.6)

# Header row
for j, (h, w) in enumerate(cat_headers):
    x = x_base_t + sum(hw[1] for hw in cat_headers[:j])
    add_shape_bg(slide, x, y_t, w, Inches(0.42), DARK_BLUE)
    add_textbox(slide, x + Inches(0.1), y_t + Inches(0.02), w - Inches(0.2), Inches(0.38),
                h, font_size=13, color=WHITE, bold=True)
y_t += Inches(0.42)

# Data rows
for i, (cat, iot, p1, p2) in enumerate(cat_rows):
    is_total = (i == len(cat_rows) - 1)
    row_h = Inches(0.42)

    # Category label
    cat_bg = DARK_BLUE if is_total else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    cat_txt = WHITE if is_total else DARK_GRAY
    x = x_base_t
    w = cat_headers[0][1]
    add_shape_bg(slide, x, y_t, w, row_h, cat_bg)
    add_textbox(slide, x + Inches(0.1), y_t + Inches(0.02), w - Inches(0.2), Inches(0.38),
                cat, font_size=12, color=cat_txt, bold=is_total)

    # IoTeye column — green tint
    x = x_base_t + cat_headers[0][1]
    w = cat_headers[1][1]
    iot_bg = RGBColor(0x1B, 0x7A, 0x43) if is_total else RGBColor(0xE8, 0xF5, 0xE9)
    iot_txt = WHITE if is_total else RGBColor(0x1B, 0x7A, 0x43)
    add_shape_bg(slide, x, y_t, w, row_h, iot_bg)
    add_textbox(slide, x + Inches(0.1), y_t + Inches(0.02), w - Inches(0.2), Inches(0.38),
                iot, font_size=13, color=iot_txt, bold=True, alignment=PP_ALIGN.CENTER)

    # P1 column
    x = x_base_t + cat_headers[0][1] + cat_headers[1][1]
    w = cat_headers[2][1]
    p1_bg = ACCENT_BLUE if is_total else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    p1_txt = WHITE if is_total else DARK_GRAY
    add_shape_bg(slide, x, y_t, w, row_h, p1_bg)
    add_textbox(slide, x + Inches(0.1), y_t + Inches(0.02), w - Inches(0.2), Inches(0.38),
                p1, font_size=12, color=p1_txt, bold=is_total, alignment=PP_ALIGN.CENTER)

    # P2 column
    x = x_base_t + cat_headers[0][1] + cat_headers[1][1] + cat_headers[2][1]
    w = cat_headers[3][1]
    p2_bg = ACCENT_BLUE if is_total else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    p2_txt = WHITE if is_total else DARK_GRAY
    add_shape_bg(slide, x, y_t, w, row_h, p2_bg)
    add_textbox(slide, x + Inches(0.1), y_t + Inches(0.02), w - Inches(0.2), Inches(0.38),
                p2, font_size=12, color=p2_txt, bold=is_total, alignment=PP_ALIGN.CENTER)

    y_t += row_h

# Right side: IoTeye Exclusive Advantages callout
add_shape_bg(slide, Inches(8.6), Inches(1.8), Inches(4.4), Inches(4.5), LIGHT_BLUE)
add_textbox(slide, Inches(8.8), Inches(1.9), Inches(4.0), Inches(0.45),
            "IoTeye-Exclusive Advantages", font_size=18, color=DARK_BLUE, bold=True)
tf = add_textbox(slide, Inches(8.8), Inches(2.4), Inches(4.0), Inches(3.7),
                 "", font_size=13, color=DARK_GRAY)
add_bullet_list(tf, [
    "BrainClaw AI \u2014 94+ tools already in production at Arc Mercer",
    "Badge ID + PIN login for frontline staff without email",
    "WCAG AAA accessibility (7:1 contrast, screen reader tested)",
    "AWS HIPAA-eligible infrastructure, 99.9% uptime SLA",
    "Severity-based SLAs (Critical: 2hr response)",
    "Source code ownership option (Option A)",
    "Multi-IDP SSO migration path (Google \u2192 Microsoft)",
    "Anomaly detection & immutable audit trails",
], font_size=12, spacing=Pt(7))

# Bottom green banner
add_shape_bg(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.7), GREEN)
add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.6),
            "\u2605  36 of 37 features \u2014 Magic Link excluded (Badge ID + PIN better serves frontline staff)  |  10 features added, all within $110K budget  \u2605",
            font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 15: Why IoTeye & Next Steps ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header_bar(slide, "Why IoTeye \u2014 Next Steps")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "Why We're the Right Partner", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5),
                 "", font_size=16, color=DARK_GRAY)
add_bullet_list(tf, [
    "5+ years in special care software \u2014 we know your workforce",
    "Arc Mercer is an existing SpringBoard & Guardian customer",
    "BrainClaw AI agents already running at Arc Mercer (37+57 tools)",
    "Badge ID + PIN login for frontline staff without email",
    "Ranked-choice committee voting \u2014 matching your current process",
    "Full engagement suite: TAC RSVP, culture groups, volunteering",
    "Caf\u00e9 cashier screen + in-person store QR redemption",
    "Dark mode for night-shift staff at group homes",
    "WCAG AAA accessibility \u2014 7:1 contrast, screen reader tested",
    "AI-assisted add-ons FREE \u2014 pay only AI LLM usage (pass-through)",
], font_size=14, color=DARK_GRAY, spacing=Pt(7))

add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
            "Next Steps", font_size=24, color=ACCENT_BLUE, bold=True)
tf = add_textbox(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(3),
                 "", font_size=16, color=DARK_GRAY)
add_bullet_list(tf, [
    "Select pricing option (A, B, or C)",
    "Schedule kickoff & discovery workshop",
    "Initiate Paycom SFTP access setup",
    "Configure Google Workspace SSO + Badge ID mapping",
    "Phase One live in 12 weeks",
], font_size=16, color=DARK_GRAY, spacing=Pt(12))

# Contact box
add_shape_bg(slide, Inches(7), Inches(4.8), Inches(5), Inches(2), LIGHT_BLUE)
add_textbox(slide, Inches(7.3), Inches(4.9), Inches(4.5), Inches(0.5),
            "Contact Us", font_size=22, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(7.3), Inches(5.4), Inches(4.5), Inches(1.2),
            "IoTeye Inc.\nBasking Ridge, New Jersey\nmingjye.sheng@ioteyeinc.com\nwww.ioteyeinc.com",
            font_size=16, color=DARK_GRAY)

# Thank you
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.8),
            "Thank You", font_size=40, color=ACCENT_BLUE, bold=True)


# Save
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "IoTeye - Arc Mercer Pitch.pptx")
output_path_tmp = os.path.join(output_dir, "IoTeye - Arc Mercer Pitch - new.pptx")
try:
    prs.save(output_path)
    print(f"Saved: {output_path}")
except PermissionError:
    prs.save(output_path_tmp)
    print(f"File locked — saved to: {output_path_tmp}")
